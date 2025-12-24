"""
Refactored indexer with callback support for progress tracking.
Supports indexing multiple directories OR single files into a specific Milvus collection.
"""
import os
import re
import ollama
from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pymilvus import connections, Collection, FieldSchema, CollectionSchema, DataType, utility
from typing import Callable, Optional, List
from services.file_types import is_supported_file
# Import BM25 Service
from services.bm25_service import BM25Service

class IndexerWithCallbacks:
    """
    Codebase indexer with callback-based progress reporting.
    Supports incremental indexing - only processes new or changed files.
    """

    def __init__(
        self,
        target_paths: List[str],
        collection_name: str,
        chunk_size: int = 1000,
        embedding_model: str = "nomic-embed-text",
        embedding_dim: Optional[int] = None,
        progress_callback: Optional[Callable] = None,
        error_callback: Optional[Callable] = None,
        indexed_files_metadata: Optional[dict] = None
    ):
        self.target_paths = target_paths
        self.raw_bucket_name = collection_name # Store raw name for MinIO
        self.collection_name = re.sub(r'[^a-zA-Z0-9_]', '_', collection_name)
        self.chunk_size = chunk_size
        self.embedding_model = embedding_model
        self.embedding_dim = embedding_dim
        self.progress_callback = progress_callback or (lambda **kwargs: None)
        self.error_callback = error_callback or (lambda **kwargs: None)
        self.should_stop = False
        self.collection = None
        # For incremental indexing: {filename: {size, last_modified}}
        self.indexed_files_metadata = indexed_files_metadata or {}
        self.new_indexed_metadata = {}  # Track newly indexed files

    def emit_progress(self, **kwargs):
        if self.progress_callback:
            self.progress_callback(**kwargs)

    def emit_error(self, **kwargs):
        if self.error_callback:
            self.error_callback(**kwargs)

    def stop(self):
        self.should_stop = True
        self.emit_progress(type='stopped', message='Indexing stopped by user')

    def validate_paths(self) -> bool:
        if not self.target_paths:
            self.emit_error(type='error', error='No paths', message='No paths specified for indexing')
            return False

        for path in self.target_paths:
            if not os.path.exists(path):
                self.emit_error(type='error', error='Invalid path', message=f'Path does not exist: {path}')
                return False
        return True

    def file_needs_indexing(self, file_path: str) -> bool:
        """Check if a file is new or has been modified since last indexing."""
        filename = os.path.basename(file_path)

        # If no previous metadata, file needs indexing
        if filename not in self.indexed_files_metadata:
            return True

        # Get current file stats
        try:
            stat = os.stat(file_path)
            current_size = stat.st_size
            current_mtime = stat.st_mtime
        except OSError:
            return True  # If we can't stat, try to index anyway

        # Compare with stored metadata
        stored = self.indexed_files_metadata[filename]
        stored_size = stored.get('size', 0)
        stored_mtime = stored.get('mtime', 0)

        # File changed if size or modification time differs
        return current_size != stored_size or current_mtime != stored_mtime

    def delete_file_from_collection(self, filename: str):
        """Delete all chunks for a specific file from the Milvus collection."""
        if self.collection is None:
            return
        try:
            # Delete by filename expression
            expr = f'filename == "{filename}"'
            self.collection.delete(expr)
            self.emit_progress(type='file_deleted', message=f'Removed old entries for {filename}')
        except Exception as e:
            self.emit_error(type='warning', message=f'Failed to delete old entries for {filename}: {str(e)}')

    def connect_to_milvus(self) -> bool:
        try:
            milvus_host = os.getenv("MILVUS_HOST", "localhost")
            milvus_port = os.getenv("MILVUS_PORT", "19530")
            connections.connect("default", host=milvus_host, port=milvus_port)
            self.emit_progress(type='milvus_connected', message=f'Connected to Milvus at {milvus_host}:{milvus_port}')
            return True
        except Exception as e:
            self.emit_error(type='error', error='Milvus connection failed', message=f'Error connecting to Milvus: {str(e)}')
            return False

    def setup_collection(self, incremental: bool = False):
        """Setup or reuse Milvus collection.

        Args:
            incremental: If True, reuse existing collection instead of dropping it.
        """
        # Auto-detect embedding dimension if not provided
        if not self.embedding_dim:
            try:
                self.emit_progress(type='detecting_dimension', message=f'Detecting embedding dimension for {self.embedding_model}...')
                test_embedding = ollama.embeddings(model=self.embedding_model, prompt="test")
                self.embedding_dim = len(test_embedding['embedding'])
                self.emit_progress(type='dimension_detected', message=f'Detected dimension: {self.embedding_dim}')
            except Exception as e:
                self.emit_error(type='error', error='Dimension detection failed', message=f'Failed to detect embedding dimension: {str(e)}')
                raise

        fields = [
            FieldSchema(name="id", dtype=DataType.INT64, is_primary=True, auto_id=True),
            FieldSchema(name="content", dtype=DataType.VARCHAR, max_length=5000),
            FieldSchema(name="filename", dtype=DataType.VARCHAR, max_length=500),
            FieldSchema(name="embedding", dtype=DataType.FLOAT_VECTOR, dim=self.embedding_dim)
        ]
        schema = CollectionSchema(fields, f"Codebase search collection for {self.collection_name}")

        if utility.has_collection(self.collection_name):
            if incremental:
                # Reuse existing collection for incremental indexing
                self.collection = Collection(self.collection_name)
                self.collection.load()
                self.emit_progress(type='collection_reused', message=f'Reusing existing collection {self.collection_name} for incremental indexing')
            else:
                self.emit_progress(type='collection_reset', message=f'Deleting existing collection {self.collection_name}')
                utility.drop_collection(self.collection_name)
                self.collection = Collection(self.collection_name, schema)
                self.emit_progress(type='collection_created', message=f'Created collection {self.collection_name} (dim={self.embedding_dim})')
        else:
            self.collection = Collection(self.collection_name, schema)
            self.emit_progress(type='collection_created', message=f'Created collection {self.collection_name} (dim={self.embedding_dim})')

    def count_files(self, incremental: bool = False) -> tuple:
        """Count total indexable files in all target paths.

        Args:
            incremental: If True, only count files that need indexing (new/changed).

        Returns:
            Tuple of (files_to_process, total_files, files_needing_indexing_paths)
        """
        total_files = 0
        files_to_process = []

        for path in self.target_paths:
            if os.path.isfile(path):
                if is_supported_file(path):
                    total_files += 1
                    if not incremental or self.file_needs_indexing(path):
                        files_to_process.append(path)
            elif os.path.isdir(path):
                for root, dirs, files in os.walk(path):
                    if '.git' in root or 'venv' in root or '__pycache__' in root or 'node_modules' in root:
                        continue
                    for file in files:
                        if is_supported_file(file):
                            total_files += 1
                            full_path = os.path.join(root, file)
                            if not incremental or self.file_needs_indexing(full_path):
                                files_to_process.append(full_path)

        return len(files_to_process), total_files, files_to_process

    def _process_single_file(self, path, files_processed, total_files, data_content, data_filename, data_embeddings, is_update: bool = False):
        """Helper to process a single file.

        Args:
            is_update: If True, this is an update to an existing file (delete old entries first)
        """
        try:
            filename = os.path.basename(path)

            # If updating an existing file, delete old entries first
            if is_update and filename in self.indexed_files_metadata:
                self.delete_file_from_collection(filename)

            self.emit_progress(
                type='file_started',
                current_file=filename,
                files_processed=files_processed,
                files_total=total_files,
                percentage=round((files_processed / total_files) * 100, 2) if total_files > 0 else 0
            )

            # Extract text based on file type
            text = ""
            file_lower = path.lower()

            try:
                if file_lower.endswith('.pdf'):
                    # Handle PDF files
                    reader = PdfReader(path)
                    for page in reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"

                elif file_lower.endswith('.docx'):
                    # Handle Word documents
                    doc = Document(path)
                    for paragraph in doc.paragraphs:
                        text += paragraph.text + "\n"
                    # Also extract text from tables
                    for table in doc.tables:
                        for row in table.rows:
                            for cell in row.cells:
                                text += cell.text + " "
                            text += "\n"

                elif file_lower.endswith('.xlsx') or file_lower.endswith('.xls'):
                    # Handle Excel spreadsheets
                    workbook = load_workbook(path, data_only=True)
                    for sheet_name in workbook.sheetnames:
                        sheet = workbook[sheet_name]
                        text += f"\n--- Sheet: {sheet_name} ---\n"
                        for row in sheet.iter_rows(values_only=True):
                            row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                            if row_text.strip():
                                text += row_text + "\n"

                elif file_lower.endswith('.pptx'):
                    # Handle PowerPoint presentations
                    prs = Presentation(path)
                    for slide_num, slide in enumerate(prs.slides, 1):
                        text += f"\n--- Slide {slide_num} ---\n"
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"

                else:
                    # Handle text files
                    with open(path, "r", encoding="utf-8", errors='ignore') as f:
                        text = f.read()

            except Exception as extract_error:
                self.emit_error(
                    type='file_error',
                    current_file=filename,
                    error=str(extract_error),
                    message=f'Error extracting {filename}: {str(extract_error)}'
                )
                return 0

            if not text.strip():
                return 0

            chunks = [text[i:i+self.chunk_size] for i in range(0, len(text), self.chunk_size)]

            for chunk in chunks:
                if self.should_stop: break

                response = ollama.embeddings(model=self.embedding_model, prompt=chunk)
                embedding = response['embedding']

                data_content.append(chunk)
                data_filename.append(filename)
                data_embeddings.append(embedding)

            # Save file metadata for incremental indexing
            try:
                stat = os.stat(path)
                self.new_indexed_metadata[filename] = {
                    'size': stat.st_size,
                    'mtime': stat.st_mtime
                }
            except OSError:
                pass

            return len(chunks)

        except Exception as e:
            self.emit_error(type='file_error', current_file=filename, error=str(e), message=f'Error processing {filename}: {str(e)}')
            return 0

    def process_files_incremental(self, files_to_process: List[str], total_to_process: int):
        """Process only the files that need indexing (new or changed).

        Args:
            files_to_process: List of file paths that need indexing
            total_to_process: Total number of files to process
        """
        data_content = []
        data_filename = []
        data_embeddings = []
        files_processed = 0
        total_chunks = 0

        self.emit_progress(
            type='started',
            files_total=total_to_process,
            files_processed=0,
            chunks_total=0,
            percentage=0.0,
            message=f'Starting incremental indexing (Chunk size: {self.chunk_size}, {total_to_process} files to process)'
        )

        for file_path in files_to_process:
            if self.should_stop:
                break

            filename = os.path.basename(file_path)
            # Check if this is an update (file exists in previous metadata)
            is_update = filename in self.indexed_files_metadata

            chunks_added = self._process_single_file(
                file_path, files_processed, total_to_process,
                data_content, data_filename, data_embeddings,
                is_update=is_update
            )
            total_chunks += chunks_added
            files_processed += 1

            action = "Updated" if is_update else "Indexed"
            self.emit_progress(
                type='file_completed',
                current_file=filename,
                files_processed=files_processed,
                files_total=total_to_process,
                chunks_total=total_chunks,
                percentage=round((files_processed / total_to_process) * 100, 2) if total_to_process > 0 else 0,
                message=f'{action} {filename}'
            )

        return [data_content, data_filename, data_embeddings]

    def get_final_indexed_metadata(self) -> dict:
        """Get the combined metadata of all indexed files (existing + newly indexed)."""
        # Start with existing metadata
        final_metadata = dict(self.indexed_files_metadata)
        # Add/update with newly indexed files
        final_metadata.update(self.new_indexed_metadata)
        return final_metadata

    def run(self):
        try:
            if not self.validate_paths(): return
            if not self.connect_to_milvus(): return

            # Determine if we should do incremental indexing
            has_existing_index = bool(self.indexed_files_metadata)

            self.emit_progress(type='counting_files', message='Analyzing files...')
            files_to_process_count, total_files, files_to_process = self.count_files(incremental=has_existing_index)

            if has_existing_index:
                self.emit_progress(
                    type='files_counted',
                    files_total=total_files,
                    files_to_process=files_to_process_count,
                    message=f'Found {total_files} files, {files_to_process_count} new/changed'
                )
            else:
                self.emit_progress(type='files_counted', files_total=total_files, message=f'Found {total_files} files')

            if files_to_process_count == 0:
                self.emit_progress(type='complete', files_total=total_files, percentage=100.0, message='All files are up to date, nothing to index')
                return

            # Setup collection (incremental if we have existing index)
            self.setup_collection(incremental=has_existing_index)

            # Process files
            my_data = self.process_files_incremental(files_to_process, files_to_process_count)

            if self.should_stop: return

            if my_data[0]:
                # 1. Insert into Milvus (Vector Search)
                self.emit_progress(type='inserting_data', message=f'Inserting {len(my_data[0])} vectors...')
                self.collection.insert(my_data)

                # Only create index if it's a new collection
                if not has_existing_index:
                    self.emit_progress(type='creating_index', message='Creating vector index...')
                    index_params = {"metric_type": "L2", "index_type": "IVF_FLAT", "params": {"nlist": 128}}
                    self.collection.create_index("embedding", index_params)

                self.collection.load()

                # 2. Build BM25 Index (Keyword Search) - always rebuild for now
                self.emit_progress(type='indexing_bm25', message='Updating keyword index (BM25)...')
                try:
                    bm25_service = BM25Service()
                    # For incremental, we need to rebuild BM25 with all content
                    # TODO: In the future, could implement incremental BM25
                    bm25_service.build_index(contents=my_data[0], filenames=my_data[1])
                    bm25_service.save_to_minio(self.raw_bucket_name)
                    self.emit_progress(type='bm25_saved', message='Keyword index saved successfully')
                except Exception as e:
                    self.emit_error(type='error', error='BM25 Error', message=f'Failed to build/save BM25 index: {e}')

                mode = "Incremental indexing" if has_existing_index else "Indexing"
                self.emit_progress(
                    type='complete',
                    files_total=total_files,
                    files_processed=files_to_process_count,
                    chunks_total=len(my_data[0]),
                    embedding_dim=self.embedding_dim,
                    percentage=100.0,
                    indexed_metadata=self.get_final_indexed_metadata(),
                    message=f'{mode} completed successfully!'
                )
            else:
                self.emit_progress(type='complete', files_total=total_files, embedding_dim=self.embedding_dim, percentage=100.0, message='No content extracted')

        except Exception as e:
            self.emit_error(type='error', error='Indexing failed', message=f'Fatal error: {str(e)}')
            raise